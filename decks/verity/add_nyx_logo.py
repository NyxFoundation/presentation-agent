#!/usr/bin/env python3
"""Overlay the Nyx Foundation wordmark on the bottom-right of every slide.

Keeps the source deck untouched and writes a new file. Run from anywhere:

    python3 decks/verity/add_nyx_logo.py

Requires: python-pptx (`pip install python-pptx`).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = Path(__file__).resolve().parent
SOURCE = HERE / "Verity.pptx"
OUTPUT = HERE / "Verity_nyx.pptx"
LOGO = HERE / "assets" / "nyx_logo.png"

# Wordmark intrinsic aspect ratio (public/images/nyx_logo.svg is 1013x221).
LOGO_W, LOGO_H = 1013, 221

# Bottom-right placement, consistent across all slides.
WIDTH = Inches(1.6)
MARGIN_RIGHT = Inches(0.3)
MARGIN_BOTTOM = Inches(0.28)


def main() -> None:
    prs = Presentation(str(SOURCE))
    height = Emu(int(WIDTH * LOGO_H / LOGO_W))
    left = prs.slide_width - MARGIN_RIGHT - WIDTH
    top = prs.slide_height - MARGIN_BOTTOM - height
    count = 0
    for slide in prs.slides:
        slide.shapes.add_picture(str(LOGO), left, top, width=WIDTH, height=height)
        count += 1
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT.name} ({count} slides)")


if __name__ == "__main__":
    main()
